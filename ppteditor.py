from pynse import *
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import aspose.slides as slides
import aspose.pydrawing as drawing
import datetime
from google.oauth2 import service_account
from google.cloud import storage
from google.oauth2 import service_account
import time
import json
from imgurpython import ImgurClient
with open("keys.json", "r") as f:
    keys = json.load(f)

client = ImgurClient(keys['imgur_client_id'], keys['imgur_client_secret'])
datalist = []
date = ''

def dataGetter():
    global datalist
    global date
    nse = Nse()
    data = pd.DataFrame()
    data = nse.fii_dii()
    datalist = []
    for i in data:
        date = data[i].keys()[0]
        print(data[i].keys()[0])
        for j in data[i]:
            # pass
            datalist.append(j)
            print(j)

    datalist.pop(0)
    datalist.pop(3)
    datalist[0],datalist[1],datalist[2] = datalist[2],datalist[0],datalist[1]
    datalist[3],datalist[4],datalist[5] = datalist[5],datalist[3],datalist[4]

def pptEditor():
    global datalist
    prs = Presentation('InstData.pptx')
    i = 0
    slides = prs.slides
    for slide in slides:
        # print(slide)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            if('Crores' in text_frame.paragraphs[0].text):
                text_frame.paragraphs[0].text = str(datalist[i]) + ' Crores'
                font = text_frame.paragraphs[0].font
                font.name = 'Montserrat SemiBold'
                font.size = Pt(200)
                font.bold = True
                if(i == 1 or i == 4 or (float(datalist[i])>0 and (i == 0 or i == 3))):
                    font.color.rgb = RGBColor(70, 202, 133)
                else:
                    font.color.rgb = RGBColor(255, 0, 0)
                i+=1
                print(text_frame.paragraphs[0].text)
    prs.save('InstData.pptx')   

def pptImg():
    pres = slides.Presentation("InstData.pptx")

    desiredX = 4000
    desiredY = 5000
    scaleX = (float)(1.0 / pres.slide_size.size.width) * desiredX
    scaleY = (float)(1.0 / pres.slide_size.size.height) * desiredY

    for index in range(pres.slides.length):
        slide = pres.slides[index]
        slide.get_thumbnail(scaleX, scaleY).save("slide_{i}.jpg".format(i = index), drawing.imaging.ImageFormat.jpeg)


scopes = ['https://www.googleapis.com/auth/cloud-platform']
credent ="argon-retina-371812-7390faca29af.json"
credential = service_account.Credentials.from_service_account_file(credent, scopes=scopes)
url = 'https://storage.googleapis.com/instppt/slide_0.jpg'

def upload_imgur():
    global url
    res = client.upload_from_path("slide_0.jpg")
    print(res)
    url= res['link']

def upload_file():
    client = storage.Client(credentials=credential, project='myproject')
    bucket = client.get_bucket('instppt')
    blob = bucket.blob('slide_0.jpg')
    blob.upload_from_filename('slide_0.jpg')
    print('File Uploaded to Bucket')

def generate_download_signed_url_v4():
    global url
    """Generates a v4 signed URL for downloading a blob.

    Note that this method requires a service account key file. You can not use
    this if you are using Application Default Credentials from Google Compute
    Engine or from the Google Cloud SDK.
    """
    bucket_name = 'instppt'
    blob_name = 'slide_0.jpg'

    storage_client = storage.Client(credentials= credential)
    bucket = storage_client.bucket(bucket_name)
    blob = bucket.blob(blob_name)

    url = blob.generate_signed_url(
        version="v4",
        # This URL is valid for 15 minutes
        expiration=datetime.timedelta(minutes=15),
        # Allow GET requests using this URL.
        method="GET",
        credentials= credential
    )
    print('Got URL - ',url)

access_token = keys['access_token']
ig_id = keys['ig_id']
creation_id = ''
def renew_token():
    global access_token
    
    r_at = requests.get(f"https://graph.facebook.com/v15.0/oauth/access_token?grant_type=fb_exchange_token&client_id=1872500803126851&client_secret=31cd69aa756dc622f32f6807df0b9538&fb_exchange_token={access_token}")
    # print(r_at.json())
    if r_at.status_code == 200:
        access_token = r_at.json()['access_token']
        print("Successfuly got Access token: ",r_at.json()['access_token'])
        return
    print('Not able to get access_token')
caption=''
def upload_img():
    global url
    global ig_id
    global access_token
    global creation_id
    r1 = requests.post(f"https://graph.facebook.com/v15.0/{ig_id}/media?image_url={url}&caption={caption}&access_token={access_token}")
    if r1.status_code != 200:
        exit(r1.json())
    print('Successfuly got Creation ID: ',r1.json()['id'])

    creation_id = r1.json()['id']

def post_publish():
    r2 = requests.post(f"https://graph.facebook.com/v15.0/{ig_id}/media_publish?creation_id={creation_id}&access_token={access_token}")
    if r2.status_code != 200:
        exit(r2.json())
    print('Post Published')
done = False
while(True):
    try:
        dataGetter()
        renew_token()
        keys['access_token'] = access_token
        with open("keys.json", "w") as f:
            json.dump(keys, f)
        if(date==str(datetime.datetime.today().strftime('%d-%b-%Y'))):
            if(done==False):
                caption ='FII-DII Data for ' + str(datetime.datetime.today().strftime('%d-%b-%Y'))
                pptEditor()
                pptImg()
                upload_imgur()
                # upload_file()
                # generate_download_signed_url_v4()  
                upload_img()
                post_publish()
                done = True
        else:
            done = False
    except:
        pass
    time.sleep(100)
